from typing import Dict, List, Union

from pptx import Presentation

TextsType = List[str]
PageType = Dict[str, Union[str, TextsType]]
ContentsType = List[PageType]


def extract_content(pptx: Presentation) -> ContentsType:
    contents: ContentsType = []

    for slide in pptx.slides:
        page: PageType = {
            'title': slide.shapes.title.text,
            'texts': [],
        }

        for shape in slide.shapes:
            if shape.text == page['title']:
                continue
            page['texts'].append(shape.text)  # type: ignore

        contents.append(page)

    return contents
